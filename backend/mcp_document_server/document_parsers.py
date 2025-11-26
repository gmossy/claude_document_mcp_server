"""
Document parsing utilities for Word, PDF, and Excel files.

This module provides functions to extract text content from various document formats
including Microsoft Word (.docx), PDF files, and Excel (.xlsx) files.
"""

import base64
import io
import mimetypes
from pathlib import Path
from typing import Optional

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
except ImportError:
    SimpleDocTemplate = None

try:
    from openpyxl import load_workbook, Workbook
except ImportError:
    load_workbook = None
    Workbook = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    Presentation = None
    Inches = None
    Pt = None
    CategoryChartData = None
    XL_CHART_TYPE = None

try:
    import magic
except ImportError:
    magic = None


class DocumentParseError(Exception):
    """Raised when document parsing fails."""
    pass


class UnsupportedFormatError(Exception):
    """Raised when document format is not supported."""
    pass


def detect_file_type(file_path: Path) -> str:
    """
    Detect the MIME type of a file.

    Args:
        file_path: Path to the file

    Returns:
        MIME type string (e.g., 'application/pdf')
    """
    # Try python-magic first (more accurate)
    if magic:
        try:
            mime = magic.Magic(mime=True)
            return mime.from_file(str(file_path))
        except Exception:
            pass

    # Fallback to mimetypes
    mime_type, _ = mimetypes.guess_type(str(file_path))
    return mime_type or "application/octet-stream"


def extract_text_from_docx(file_path: Path) -> tuple[str, dict]:
    """
    Extract text content from a Word document.

    Args:
        file_path: Path to the .docx file

    Returns:
        Tuple of (extracted_text, metadata_dict)

    Raises:
        DocumentParseError: If parsing fails
        UnsupportedFormatError: If python-docx is not installed
    """
    if DocxDocument is None:
        raise UnsupportedFormatError(
            "python-docx is not installed. Install with: pip install python-docx"
        )

    try:
        doc = DocxDocument(file_path)

        # Extract all paragraph text
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        text = "\n\n".join(paragraphs)

        # Extract metadata
        metadata = {
            "format": "docx",
            "paragraph_count": len(doc.paragraphs),
            "has_tables": len(doc.tables) > 0,
            "table_count": len(doc.tables),
        }

        # Try to get core properties
        try:
            core_props = doc.core_properties
            if core_props.author:
                metadata["author"] = core_props.author
            if core_props.title:
                metadata["original_title"] = core_props.title
            if core_props.subject:
                metadata["subject"] = core_props.subject
            if core_props.created:
                metadata["created"] = core_props.created.isoformat()
            if core_props.modified:
                metadata["modified"] = core_props.modified.isoformat()
        except Exception:
            pass

        return text, metadata

    except Exception as e:
        raise DocumentParseError(f"Failed to parse Word document: {str(e)}")


def extract_text_from_pdf(file_path: Path) -> tuple[str, dict]:
    """
    Extract text content from a PDF file.

    Args:
        file_path: Path to the .pdf file

    Returns:
        Tuple of (extracted_text, metadata_dict)

    Raises:
        DocumentParseError: If parsing fails
        UnsupportedFormatError: If pypdf is not installed
    """
    if PdfReader is None:
        raise UnsupportedFormatError(
            "pypdf is not installed. Install with: pip install pypdf"
        )

    try:
        reader = PdfReader(file_path)

        # Extract text from all pages
        pages_text = []
        for page_num, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text()
                if text.strip():
                    pages_text.append(f"--- Page {page_num} ---\n{text}")
            except Exception:
                pages_text.append(f"--- Page {page_num} ---\n[Could not extract text]")

        text = "\n\n".join(pages_text)

        # Extract metadata
        metadata = {
            "format": "pdf",
            "page_count": len(reader.pages),
        }

        # Try to get PDF metadata
        try:
            if reader.metadata:
                if reader.metadata.author:
                    metadata["author"] = reader.metadata.author
                if reader.metadata.title:
                    metadata["original_title"] = reader.metadata.title
                if reader.metadata.subject:
                    metadata["subject"] = reader.metadata.subject
                if reader.metadata.creator:
                    metadata["creator"] = reader.metadata.creator
                if reader.metadata.producer:
                    metadata["producer"] = reader.metadata.producer
                if reader.metadata.creation_date:
                    metadata["created"] = str(reader.metadata.creation_date)
        except Exception:
            pass

        return text, metadata

    except Exception as e:
        raise DocumentParseError(f"Failed to parse PDF document: {str(e)}")


def extract_text_from_excel(file_path: Path) -> tuple[str, dict]:
    """
    Extract text content from an Excel file.

    Args:
        file_path: Path to the .xlsx file

    Returns:
        Tuple of (extracted_text, metadata_dict)

    Raises:
        DocumentParseError: If parsing fails
        UnsupportedFormatError: If openpyxl is not installed
    """
    if load_workbook is None:
        raise UnsupportedFormatError(
            "openpyxl is not installed. Install with: pip install openpyxl"
        )

    try:
        workbook = load_workbook(file_path, data_only=True)

        # Extract text from all sheets
        sheets_text = []
        total_rows = 0
        total_cells = 0

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_data = []

            # Get the used range
            if sheet.max_row > 0:
                sheet_data.append(f"=== Sheet: {sheet_name} ===\n")

                for row in sheet.iter_rows(values_only=True):
                    # Filter out completely empty rows
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(val.strip() for val in row_values):
                        # Join cells with tabs to preserve column structure
                        sheet_data.append("\t".join(row_values))
                        total_rows += 1
                        total_cells += len([v for v in row_values if v.strip()])

                if len(sheet_data) > 1:  # More than just the header
                    sheets_text.append("\n".join(sheet_data))

        text = "\n\n".join(sheets_text)

        # Extract metadata
        metadata = {
            "format": "xlsx",
            "sheet_count": len(workbook.sheetnames),
            "sheet_names": workbook.sheetnames,
            "total_rows": total_rows,
            "total_cells": total_cells,
        }

        # Try to get Excel properties
        try:
            props = workbook.properties
            if props:
                if props.creator:
                    metadata["author"] = props.creator
                if props.title:
                    metadata["original_title"] = props.title
                if props.subject:
                    metadata["subject"] = props.subject
                if props.created:
                    metadata["created"] = str(props.created)
                if props.modified:
                    metadata["modified"] = str(props.modified)
        except Exception:
            pass

        workbook.close()
        return text, metadata

    except Exception as e:
        raise DocumentParseError(f"Failed to parse Excel document: {str(e)}")


def extract_text_from_powerpoint(file_path: Path) -> tuple[str, dict]:
    """
    Extract text content from a PowerPoint file.

    Args:
        file_path: Path to the .pptx file

    Returns:
        Tuple of (extracted_text, metadata_dict)

    Raises:
        DocumentParseError: If parsing fails
        UnsupportedFormatError: If python-pptx is not installed
    """
    if Presentation is None:
        raise UnsupportedFormatError(
            "python-pptx is not installed. Install with: pip install python-pptx"
        )

    try:
        prs = Presentation(file_path)

        # Extract text from all slides
        slides_text = []
        total_shapes = 0
        total_text_boxes = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            slide_content.append(f"=== Slide {slide_num} ===\n")

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                total_shapes += 1

                # Check if shape has text
                if hasattr(shape, "text") and shape.text.strip():
                    total_text_boxes += 1
                    slide_content.append(shape.text)

                # Check for tables
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        if any(row_data):  # Only add non-empty rows
                            table_data.append("\t".join(row_data))
                    if table_data:
                        slide_content.append("\n[Table]")
                        slide_content.extend(table_data)

            if len(slide_content) > 1:  # More than just the slide header
                slides_text.append("\n".join(slide_content))

        text = "\n\n".join(slides_text)

        # Extract metadata
        metadata = {
            "format": "pptx",
            "slide_count": len(prs.slides),
            "total_shapes": total_shapes,
            "text_boxes": total_text_boxes,
        }

        # Try to get PowerPoint properties
        try:
            core_props = prs.core_properties
            if core_props:
                if core_props.author:
                    metadata["author"] = core_props.author
                if core_props.title:
                    metadata["original_title"] = core_props.title
                if core_props.subject:
                    metadata["subject"] = core_props.subject
                if core_props.created:
                    metadata["created"] = str(core_props.created)
                if core_props.modified:
                    metadata["modified"] = str(core_props.modified)
                if core_props.last_modified_by:
                    metadata["last_modified_by"] = core_props.last_modified_by
        except Exception:
            pass

        return text, metadata

    except Exception as e:
        raise DocumentParseError(f"Failed to parse PowerPoint document: {str(e)}")


def extract_text_from_file(file_path: Path) -> tuple[str, dict]:
    """
    Auto-detect file type and extract text content.

    Args:
        file_path: Path to the document file

    Returns:
        Tuple of (extracted_text, metadata_dict)

    Raises:
        UnsupportedFormatError: If file format is not supported
        DocumentParseError: If parsing fails
    """
    if not file_path.exists():
        raise DocumentParseError(f"File not found: {file_path}")

    # Detect file type
    mime_type = detect_file_type(file_path)
    extension = file_path.suffix.lower()

    # Route to appropriate parser
    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or extension == ".docx":
        return extract_text_from_docx(file_path)
    elif mime_type == "application/pdf" or extension == ".pdf":
        return extract_text_from_pdf(file_path)
    elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or extension == ".xlsx":
        return extract_text_from_excel(file_path)
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or extension == ".pptx":
        return extract_text_from_powerpoint(file_path)
    elif mime_type.startswith("text/") or extension in [".txt", ".md", ".markdown"]:
        # Plain text files
        text = file_path.read_text(encoding="utf-8")
        metadata = {"format": "text", "mime_type": mime_type}
        return text, metadata
    else:
        raise UnsupportedFormatError(
            f"Unsupported file format: {mime_type} ({extension}). "
            f"Supported formats: .docx, .pdf, .xlsx, .pptx, .txt, .md"
        )


def create_pdf_from_text(text: str, output_path: Path, title: str = "Document") -> None:
    """
    Create a PDF file from text content.

    Args:
        text: Text content to write to PDF
        output_path: Path where PDF should be saved
        title: Document title

    Raises:
        UnsupportedFormatError: If reportlab is not installed
        DocumentParseError: If PDF creation fails
    """
    if SimpleDocTemplate is None:
        raise UnsupportedFormatError(
            "reportlab is not installed. Install with: pip install reportlab"
        )

    try:
        doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Add title
        title_para = Paragraph(title, styles['Title'])
        story.append(title_para)
        story.append(Spacer(1, 0.2 * inch))

        # Add content paragraphs
        for paragraph in text.split('\n\n'):
            if paragraph.strip():
                para = Paragraph(paragraph.replace('\n', '<br/>'), styles['BodyText'])
                story.append(para)
                story.append(Spacer(1, 0.1 * inch))

        doc.build(story)

    except Exception as e:
        raise DocumentParseError(f"Failed to create PDF: {str(e)}")


def create_docx_from_text(text: str, output_path: Path, title: str = "Document") -> None:
    """
    Create a Word document from text content.

    Args:
        text: Text content to write to document
        output_path: Path where .docx should be saved
        title: Document title

    Raises:
        UnsupportedFormatError: If python-docx is not installed
        DocumentParseError: If document creation fails
    """
    if DocxDocument is None:
        raise UnsupportedFormatError(
            "python-docx is not installed. Install with: pip install python-docx"
        )

    try:
        doc = DocxDocument()

        # Add title
        doc.add_heading(title, 0)

        # Add content paragraphs
        for paragraph in text.split('\n\n'):
            if paragraph.strip():
                doc.add_paragraph(paragraph)

        doc.save(str(output_path))

    except Exception as e:
        raise DocumentParseError(f"Failed to create Word document: {str(e)}")


def create_excel_from_data(data: list[list[str]], output_path: Path, sheet_name: str = "Sheet1", title: str = "Document") -> None:
    """
    Create an Excel file from tabular data.

    Args:
        data: List of rows, where each row is a list of cell values
        output_path: Path where Excel file should be saved
        sheet_name: Name of the worksheet
        title: Document title (used in properties)

    Raises:
        UnsupportedFormatError: If openpyxl is not installed
        DocumentParseError: If Excel creation fails
    """
    if Workbook is None:
        raise UnsupportedFormatError(
            "openpyxl is not installed. Install with: pip install openpyxl"
        )

    try:
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = sheet_name

        # Write data to sheet
        for row_data in data:
            sheet.append(row_data)

        # Set document properties
        workbook.properties.title = title

        # Auto-adjust column widths (basic)
        for column in sheet.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)  # Cap at 50
            sheet.column_dimensions[column_letter].width = adjusted_width

        workbook.save(str(output_path))
        workbook.close()

    except Exception as e:
        raise DocumentParseError(f"Failed to create Excel document: {str(e)}")


def create_powerpoint_from_slides(slides_data: list[dict], output_path: Path, title: str = "Presentation", template_path: Optional[Path] = None) -> None:
    """
    Create a PowerPoint file from slide data, optionally using a template.

    Args:
        slides_data: List of slide dictionaries, each containing:
            - title: Slide title (optional)
            - content: List of bullet points or text blocks
            - layout: Layout type (default: "title_and_content")
            - layout_index: Specific layout index from template (optional)
            - chart: Chart data (optional) with:
                - type: Chart type ("bar", "column", "line", "pie")
                - categories: List of category labels
                - series: List of dicts with 'name' and 'values'
        output_path: Path where PowerPoint file should be saved
        title: Presentation title (used in properties)
        template_path: Path to existing .pptx template file (optional)
                       If provided, uses template's theme, colors, fonts, and layouts

    Raises:
        UnsupportedFormatError: If python-pptx is not installed
        DocumentParseError: If PowerPoint creation fails
    """
    if Presentation is None or Inches is None or Pt is None:
        raise UnsupportedFormatError(
            "python-pptx is not installed. Install with: pip install python-pptx"
        )

    try:
        # Load template if provided, otherwise create blank presentation
        if template_path and template_path.exists():
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()
            prs.slide_width = Inches(10)  # Standard 16:9
            prs.slide_height = Inches(7.5)

        # Set presentation properties
        prs.core_properties.title = title

        for slide_data in slides_data:
            # Determine layout
            layout_type = slide_data.get("layout", "title_and_content")

            # If layout_index is specified, use that exact layout from template
            if "layout_index" in slide_data:
                layout_index = slide_data["layout_index"]
                if 0 <= layout_index < len(prs.slide_layouts):
                    slide_layout = prs.slide_layouts[layout_index]
                else:
                    slide_layout = prs.slide_layouts[1]  # Default to title and content
            else:
                # Use layout type string (for backward compatibility)
                if layout_type == "title":
                    slide_layout = prs.slide_layouts[0]  # Title slide
                elif layout_type == "blank":
                    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[1]
                else:
                    slide_layout = prs.slide_layouts[1]  # Title and content

            slide = prs.slides.add_slide(slide_layout)

            # Add title if provided
            if "title" in slide_data and slide_data["title"]:
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data["title"]

            # Add content
            if "content" in slide_data:
                content = slide_data["content"]

                # For title and content layout, use the text placeholder
                if layout_type == "title_and_content" and len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.clear()  # Clear default text

                    if isinstance(content, list):
                        # Add bullet points
                        for i, item in enumerate(content):
                            if i == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            p.text = str(item)
                            p.level = 0
                    else:
                        # Add as single text block
                        text_frame.text = str(content)

                # For blank layout, add text box
                elif layout_type == "blank":
                    left = Inches(1)
                    top = Inches(1.5)
                    width = Inches(8)
                    height = Inches(5)

                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame

                    if isinstance(content, list):
                        text_frame.text = "\n".join(str(item) for item in content)
                    else:
                        text_frame.text = str(content)

            # Add chart if provided
            if "chart" in slide_data and CategoryChartData is not None and XL_CHART_TYPE is not None:
                chart_info = slide_data["chart"]
                chart_type_str = chart_info.get("type", "column").lower()

                # Map chart type string to enum
                chart_type_map = {
                    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                    "line": XL_CHART_TYPE.LINE,
                    "pie": XL_CHART_TYPE.PIE,
                }
                chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

                # Create chart data
                chart_data = CategoryChartData()
                chart_data.categories = chart_info.get("categories", [])

                # Add series
                for series_info in chart_info.get("series", []):
                    chart_data.add_series(
                        series_info.get("name", "Series"),
                        series_info.get("values", [])
                    )

                # Add chart to slide
                x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(5)
                chart = slide.shapes.add_chart(
                    chart_type, x, y, cx, cy, chart_data
                ).chart

                # Optional: Set chart title if provided
                if "chart_title" in chart_info:
                    chart.has_title = True
                    chart.chart_title.text_frame.text = chart_info["chart_title"]

        prs.save(str(output_path))

    except Exception as e:
        raise DocumentParseError(f"Failed to create PowerPoint document: {str(e)}")


def get_powerpoint_template_layouts(template_path: Path) -> list[dict]:
    """
    Get information about available layouts in a PowerPoint template.

    Args:
        template_path: Path to the .pptx template file

    Returns:
        List of layout dictionaries with 'index', 'name', and 'placeholder_count'

    Raises:
        UnsupportedFormatError: If python-pptx is not installed
        DocumentParseError: If template cannot be read
    """
    if Presentation is None:
        raise UnsupportedFormatError(
            "python-pptx is not installed. Install with: pip install python-pptx"
        )

    try:
        prs = Presentation(str(template_path))
        layouts = []

        for idx, layout in enumerate(prs.slide_layouts):
            layout_info = {
                "index": idx,
                "name": layout.name,
                "placeholder_count": len(layout.placeholders),
                "placeholders": []
            }

            # Get placeholder information
            for placeholder in layout.placeholders:
                placeholder_info = {
                    "index": placeholder.placeholder_format.idx,
                    "type": str(placeholder.placeholder_format.type),
                    "name": placeholder.name
                }
                layout_info["placeholders"].append(placeholder_info)

            layouts.append(layout_info)

        return layouts

    except Exception as e:
        raise DocumentParseError(f"Failed to read PowerPoint template: {str(e)}")


def encode_file_to_base64(file_path: Path) -> str:
    """
    Encode a file to base64 string.

    Args:
        file_path: Path to file

    Returns:
        Base64 encoded string
    """
    with open(file_path, 'rb') as f:
        return base64.b64encode(f.read()).decode('utf-8')


def decode_base64_to_file(base64_str: str, output_path: Path) -> None:
    """
    Decode base64 string and write to file.

    Args:
        base64_str: Base64 encoded file content
        output_path: Path where file should be saved
    """
    file_data = base64.b64decode(base64_str)
    with open(output_path, 'wb') as f:
        f.write(file_data)
